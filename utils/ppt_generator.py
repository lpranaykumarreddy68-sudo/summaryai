"""
PPT Slide Generator — Produces a themed PowerPoint deck via LLM + python-pptx.
Includes high-quality relevant downloaded images for each slide theme from curated CDN assets.
"""

import io
import json
import re
import requests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from langchain_core.messages import HumanMessage


SLIDES_PROMPT = """You are an expert presentation designer. Analyze the following document and 
extract exactly 4 key themes. For each theme, create a slide outline.

Return your response as a valid JSON array with exactly 4 objects. Each object must have:
- "title": A concise slide title (max 8 words)
- "bullets": An array of 3-5 concise bullet point strings
- "keyword": A single search keyword or short phrase (1-2 words) that describes the visual theme of the slide (e.g., "robotics", "brain", "calculation", "history", "graph", "business", "coding", "physics", "drone")

Example format:
[
  {{"title": "Introduction to Topic", "bullets": ["Point 1", "Point 2", "Point 3"], "keyword": "drone"}},
  {{"title": "Key Concepts", "bullets": ["Concept A", "Concept B", "Concept C"], "keyword": "autonomous"}}
]

IMPORTANT: Return ONLY the JSON array. No markdown, no code fences, no extra text.

Document text:
{text}"""


def _parse_slides_json(raw: str) -> list:
    """Parse slides JSON from LLM response, handling common formatting issues."""
    # Strip markdown code fences if present
    cleaned = raw.strip()
    cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
    cleaned = re.sub(r"\s*```$", "", cleaned)
    cleaned = cleaned.strip()

    try:
        data = json.loads(cleaned)
    except json.JSONDecodeError:
        # Try to find JSON array in the response
        match = re.search(r"\[.*\]", cleaned, re.DOTALL)
        if match:
            data = json.loads(match.group())
        else:
            raise ValueError("Could not parse slide data from LLM response.")

    if not isinstance(data, list) or len(data) == 0:
        raise ValueError("Expected a non-empty JSON array of slide objects.")

    return data


def get_exact_image_url(keyword: str) -> str:
    """
    Return a curated, high-quality, professional Unsplash image URL for a given keyword.
    These are verified, premium, high-resolution images that match the autonomous nav/robotics themes perfectly.
    """
    keyword_lower = keyword.lower().strip()
    
    # Hand-picked premium high-quality Unsplash image assets matching common themes:
    mappings = {
        "drone": "https://images.unsplash.com/photo-1508614589041-895b88991e3e?auto=format&fit=crop&w=600&q=80",
        "uav": "https://images.unsplash.com/photo-1508614589041-895b88991e3e?auto=format&fit=crop&w=600&q=80",
        "autonomous": "https://images.unsplash.com/photo-1549317661-bd32c8ce0db2?auto=format&fit=crop&w=600&q=80", # Self-driving car interior
        "self-driving": "https://images.unsplash.com/photo-1549317661-bd32c8ce0db2?auto=format&fit=crop&w=600&q=80",
        "navigation": "https://images.unsplash.com/photo-1524661135-423995f22d0b?auto=format&fit=crop&w=600&q=80", # Radar/GPS Map
        "map": "https://images.unsplash.com/photo-1524661135-423995f22d0b?auto=format&fit=crop&w=600&q=80",
        "robot": "https://images.unsplash.com/photo-1485827404703-89b55fcc595e?auto=format&fit=crop&w=600&q=80", # Futuristic robot
        "robotics": "https://images.unsplash.com/photo-1485827404703-89b55fcc595e?auto=format&fit=crop&w=600&q=80",
        "sensor": "https://images.unsplash.com/photo-1518770660439-4636190af475?auto=format&fit=crop&w=600&q=80", # Microchip/Circuit
        "microchip": "https://images.unsplash.com/photo-1518770660439-4636190af475?auto=format&fit=crop&w=600&q=80",
        "circuit": "https://images.unsplash.com/photo-1518770660439-4636190af475?auto=format&fit=crop&w=600&q=80",
        "technology": "https://images.unsplash.com/photo-1451187580459-43490279c0fa?auto=format&fit=crop&w=600&q=80", # Connected blue network
        "innovation": "https://images.unsplash.com/photo-1451187580459-43490279c0fa?auto=format&fit=crop&w=600&q=80",
        "hub": "https://images.unsplash.com/photo-1486406146926-c627a92ad1ab?auto=format&fit=crop&w=600&q=80", # Corporate tech building
        "iith": "https://images.unsplash.com/photo-1562774053-701939374585?auto=format&fit=crop&w=600&q=80", # Indian college campus
        "iit": "https://images.unsplash.com/photo-1562774053-701939374585?auto=format&fit=crop&w=600&q=80",
        "india": "https://images.unsplash.com/photo-1524492412937-b28074a5d7da?auto=format&fit=crop&w=600&q=80", # Taj Mahal/India
        "campus": "https://images.unsplash.com/photo-1562774053-701939374585?auto=format&fit=crop&w=600&q=80",
        "vehicle": "https://images.unsplash.com/photo-1549317661-bd32c8ce0db2?auto=format&fit=crop&w=600&q=80",
        "car": "https://images.unsplash.com/photo-1549317661-bd32c8ce0db2?auto=format&fit=crop&w=600&q=80",
        "testing": "https://images.unsplash.com/photo-1581091226825-a6a2a5aee158?auto=format&fit=crop&w=600&q=80", # Engineer testing
        "testbed": "https://images.unsplash.com/photo-1581091226825-a6a2a5aee158?auto=format&fit=crop&w=600&q=80",
        "research": "https://images.unsplash.com/photo-1507679799987-c73779587ccf?auto=format&fit=crop&w=600&q=80", # Businessman/Researcher in suits
        "funding": "https://images.unsplash.com/photo-1559526324-4b87b5e36e44?auto=format&fit=crop&w=600&q=80", # Budget chart
        "education": "https://images.unsplash.com/photo-1497633762265-9d179a990aa6?auto=format&fit=crop&w=600&q=80", # Books
        "learning": "https://images.unsplash.com/photo-1497633762265-9d179a990aa6?auto=format&fit=crop&w=600&q=80",
    }

    # Match exact keywords or substrings
    for key, url in mappings.items():
        if key in keyword_lower:
            return url

    # Default fallbacks based on visual concept:
    if any(k in keyword_lower for k in ["computer", "code", "software", "ai", "artificial"]):
        return "https://images.unsplash.com/photo-1517694712202-14dd9538aa97?auto=format&fit=crop&w=600&q=80" # Mac coding screen
    if any(k in keyword_lower for k in ["space", "sky", "aerial"]):
        return "https://images.unsplash.com/photo-1451187580459-43490279c0fa?auto=format&fit=crop&w=600&q=80"
    if any(k in keyword_lower for k in ["chart", "graph", "business", "data"]):
        return "https://images.unsplash.com/photo-1551836022-d5d88e9218df?auto=format&fit=crop&w=600&q=80" # Analytics screen

    # Generic high-quality tech fallback:
    return "https://images.unsplash.com/photo-1451187580459-43490279c0fa?auto=format&fit=crop&w=600&q=80"


def _add_styled_slide(prs: Presentation, title: str, bullets: list, keyword: str, slide_num: int):
    """Add a professionally styled slide to the presentation with an optional downloaded image."""
    slide_layout = prs.slide_layouts[6]  # Blank layout for full control
    slide = prs.slides.add_slide(slide_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # --- Accent bar on left ---
    accent_colors = [
        RGBColor(59, 130, 246),   # Blue
        RGBColor(16, 185, 129),   # Green
        RGBColor(139, 92, 246),   # Purple
        RGBColor(245, 158, 11),   # Amber
    ]
    color = accent_colors[slide_num % len(accent_colors)]

    accent = slide.shapes.add_shape(
        1, Emu(0), Emu(0), Inches(0.15), slide_height  # Rectangle
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

    # --- Slide number badge ---
    badge = slide.shapes.add_shape(
        1, Inches(0.5), Inches(0.4), Inches(0.5), Inches(0.5)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf_badge = badge.text_frame
    tf_badge.text = str(slide_num + 1)
    tf_badge.paragraphs[0].font.size = Pt(18)
    tf_badge.paragraphs[0].font.bold = True
    tf_badge.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf_badge.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- Title ---
    title_box = slide.shapes.add_textbox(Inches(1.3), Inches(0.35), Inches(7.5), Inches(0.7))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p = tf_title.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 41, 59)

    # --- Divider line ---
    line = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.15), Inches(8.5), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(226, 232, 240)
    line.line.fill.background()

    # --- Download image related to topic from Unsplash CDN ---
    img_bytes = None
    if keyword and keyword.strip():
        try:
            url = get_exact_image_url(keyword)
            # Fetch with a strict timeout to avoid slowing down presentation generation
            res = requests.get(url, timeout=7)
            if res.status_code == 200 and len(res.content) > 1000:
                img_bytes = res.content
        except Exception:
            pass  # Fallback to text-only if download fails or times out

    # --- Layout dimensions (text and image) ---
    if img_bytes:
        box_width = Inches(4.7)
        img_left = Inches(5.5)
        img_top = Inches(1.7)
        img_width = Inches(3.8)
        img_height = Inches(4.5)
    else:
        box_width = Inches(8.0)

    # --- Bullet points ---
    bullet_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), box_width, Inches(4.5))
    tf_bullets = bullet_box.text_frame
    tf_bullets.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf_bullets.paragraphs[0]
        else:
            p = tf_bullets.add_paragraph()

        p.text = f"▸  {bullet}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_after = Pt(12)
        p.level = 0

    # --- Add image if successfully fetched ---
    if img_bytes:
        try:
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                img_left,
                img_top,
                width=img_width,
                height=img_height
            )
        except Exception:
            pass


def generate_ppt_slides(llm, text: str) -> bytes:
    """
    Generate a PowerPoint presentation with 4 themed slides and relevant visuals.

    Args:
        llm: A LangChain chat model.
        text: The full document text.

    Returns:
        PPTX file contents as bytes.
    """
    truncated = text[:15000] if len(text) > 15000 else text
    prompt = SLIDES_PROMPT.format(text=truncated)
    response = llm.invoke([HumanMessage(content=prompt)])
    slides_data = _parse_slides_json(response.content)

    # Build the presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Title slide ---
    title_slide_layout = prs.slide_layouts[6]  # Blank
    title_slide = prs.slides.add_slide(title_slide_layout)

    # Background accent
    bg_shape = title_slide.shapes.add_shape(
        1, Emu(0), Emu(0), prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(30, 41, 59)
    bg_shape.line.fill.background()

    # Title text
    tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "📚 AI Study Companion"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    sub = tf.add_paragraph()
    sub.text = "Key Themes & Visual Study Guide"
    sub.font.size = Pt(18)
    sub.font.color.rgb = RGBColor(148, 163, 184)
    sub.alignment = PP_ALIGN.CENTER

    # --- Content slides ---
    for i, slide_data in enumerate(slides_data[:4]):
        title = slide_data.get("title", f"Theme {i + 1}")
        bullets = slide_data.get("bullets", ["No content generated."])
        keyword = slide_data.get("keyword", "education")
        _add_styled_slide(prs, title, bullets, keyword, i)

    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()

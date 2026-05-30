"""
Generators — Houses background task functions for asset generation.
Runs PDF, PPTX, MP3, and Flashcard generation asynchronously via FastAPI BackgroundTasks,
updating the SQLite database with progress, URLs, and errors.
"""

import io
import os
import json
import re
import requests
from datetime import datetime
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from gtts import gTTS
from langchain_core.messages import HumanMessage

from database import SessionLocal, GenerationTask, Flashcard
from utils.llm_provider import get_llm

# Directory for storing generated files
STATIC_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "static")
EXPORTS_DIR = os.path.join(STATIC_DIR, "exports")
os.makedirs(EXPORTS_DIR, exist_ok=True)

# ── Curated High-Quality Visual CDN Mapping ──────────────────────────────────
UNSPLASH_MAPPINGS = {
    "drone": "https://images.unsplash.com/photo-1508614589041-895b88991e3e?auto=format&fit=crop&w=600&q=80",
    "uav": "https://images.unsplash.com/photo-1508614589041-895b88991e3e?auto=format&fit=crop&w=600&q=80",
    "autonomous": "https://images.unsplash.com/photo-1549317661-bd32c8ce0db2?auto=format&fit=crop&w=600&q=80",
    "self-driving": "https://images.unsplash.com/photo-1549317661-bd32c8ce0db2?auto=format&fit=crop&w=600&q=80",
    "navigation": "https://images.unsplash.com/photo-1524661135-423995f22d0b?auto=format&fit=crop&w=600&q=80",
    "map": "https://images.unsplash.com/photo-1524661135-423995f22d0b?auto=format&fit=crop&w=600&q=80",
    "robot": "https://images.unsplash.com/photo-1485827404703-89b55fcc595e?auto=format&fit=crop&w=600&q=80",
    "robotics": "https://images.unsplash.com/photo-1485827404703-89b55fcc595e?auto=format&fit=crop&w=600&q=80",
    "sensor": "https://images.unsplash.com/photo-1518770660439-4636190af475?auto=format&fit=crop&w=600&q=80",
    "microchip": "https://images.unsplash.com/photo-1518770660439-4636190af475?auto=format&fit=crop&w=600&q=80",
    "circuit": "https://images.unsplash.com/photo-1518770660439-4636190af475?auto=format&fit=crop&w=600&q=80",
    "technology": "https://images.unsplash.com/photo-1451187580459-43490279c0fa?auto=format&fit=crop&w=600&q=80",
    "innovation": "https://images.unsplash.com/photo-1451187580459-43490279c0fa?auto=format&fit=crop&w=600&q=80",
    "hub": "https://images.unsplash.com/photo-1486406146926-c627a92ad1ab?auto=format&fit=crop&w=600&q=80",
    "iith": "https://images.unsplash.com/photo-1562774053-701939374585?auto=format&fit=crop&w=600&q=80",
    "iit": "https://images.unsplash.com/photo-1562774053-701939374585?auto=format&fit=crop&w=600&q=80",
    "india": "https://images.unsplash.com/photo-1524492412937-b28074a5d7da?auto=format&fit=crop&w=600&q=80",
    "campus": "https://images.unsplash.com/photo-1562774053-701939374585?auto=format&fit=crop&w=600&q=80",
    "vehicle": "https://images.unsplash.com/photo-1549317661-bd32c8ce0db2?auto=format&fit=crop&w=600&q=80",
    "car": "https://images.unsplash.com/photo-1549317661-bd32c8ce0db2?auto=format&fit=crop&w=600&q=80",
    "testing": "https://images.unsplash.com/photo-1581091226825-a6a2a5aee158?auto=format&fit=crop&w=600&q=80",
    "testbed": "https://images.unsplash.com/photo-1581091226825-a6a2a5aee158?auto=format&fit=crop&w=600&q=80",
    "research": "https://images.unsplash.com/photo-1507679799987-c73779587ccf?auto=format&fit=crop&w=600&q=80",
    "funding": "https://images.unsplash.com/photo-1559526324-4b87b5e36e44?auto=format&fit=crop&w=600&q=80",
    "education": "https://images.unsplash.com/photo-1497633762265-9d179a990aa6?auto=format&fit=crop&w=600&q=80",
    "learning": "https://images.unsplash.com/photo-1497633762265-9d179a990aa6?auto=format&fit=crop&w=600&q=80",
}


def _get_exact_image_url(keyword: str) -> str:
    """Return a curated Unsplash image CDN URL for a given keyword."""
    kw_lower = keyword.lower().strip()
    for key, url in UNSPLASH_MAPPINGS.items():
        if key in kw_lower:
            return url

    # Default visual fallbacks
    if any(k in kw_lower for k in ["computer", "code", "software", "ai", "artificial"]):
        return "https://images.unsplash.com/photo-1517694712202-14dd9538aa97?auto=format&fit=crop&w=600&q=80"
    if any(k in kw_lower for k in ["space", "sky", "aerial"]):
        return "https://images.unsplash.com/photo-1451187580459-43490279c0fa?auto=format&fit=crop&w=600&q=80"
    if any(k in kw_lower for k in ["chart", "graph", "business", "data"]):
        return "https://images.unsplash.com/photo-1551836022-d5d88e9218df?auto=format&fit=crop&w=600&q=80"
    
    return "https://images.unsplash.com/photo-1451187580459-43490279c0fa?auto=format&fit=crop&w=600&q=80"


def _sanitize_latin1(text: str) -> str:
    """Sanitize text to support only FPDF Helvetica Latin-1 encoding characters."""
    replacements = {
        "\u2022": "-",
        "\u2013": "-",
        "\u2014": "--",
        "\u2018": "'",
        "\u2019": "'",
        "\u201c": '"',
        "\u201d": '"',
        "\u2026": "...",
        "\u00ae": "(R)",
        "\u00ad": "",
        "\u2122": "TM",
    }
    for uni, ascii_val in replacements.items():
        text = text.replace(uni, ascii_val)
    return text.encode("latin-1", errors="replace").decode("latin-1")


def _sanitize_error_message(e: Exception) -> str:
    """Convert low-level API rate-limit errors to clear, actionable instructions."""
    err_str = str(e)
    if "429" in err_str or "RESOURCE_EXHAUSTED" in err_str:
        return "Gemini API Quota Exceeded: The configured Google API key has exceeded the daily free tier limit of 20 requests. Please configure an alternative key in your .env file or retry tomorrow."
    return err_str


# ── Background Task: PDF Summary Generation ─────────────────────────────────
def generate_pdf_task(raw_text: str, task_id: str):
    db = SessionLocal()
    task = db.query(GenerationTask).filter(GenerationTask.task_id == task_id).first()
    if not task:
        db.close()
        return

    try:
        # Load API keys from environment
        provider = os.getenv("AI_PROVIDER", "gemini").lower().strip()
        api_key = os.getenv("GOOGLE_API_KEY" if provider == "gemini" else "OPENAI_API_KEY", "")
        llm = get_llm(provider, api_key)

        prompt = """You are an expert study assistant. Read the following document text and produce 
a comprehensive, well-structured summary suitable for student revision.

Requirements:
- Use clear section headings (prefixed with ##).
- Under each heading, use concise bullet points (prefixed with •).
- Cover ALL core concepts, definitions, and key takeaways.
- Keep language simple and student-friendly.

Document text:
{text}

Structured Summary:"""

        truncated = raw_text[:15000] if len(raw_text) > 15000 else raw_text
        response = llm.invoke([HumanMessage(content=prompt.format(text=truncated))])
        summary = _sanitize_latin1(response.content)

        # Build PDF
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=20)
        pdf.add_page()

        pdf.set_font("Helvetica", "B", 18)
        pdf.set_text_color(30, 58, 138)
        pdf.cell(0, 12, "AI Study Companion", new_x="LMARGIN", new_y="NEXT", align="C")

        pdf.set_font("Helvetica", "", 11)
        pdf.set_text_color(100, 100, 100)
        pdf.cell(0, 8, "Document Summary", new_x="LMARGIN", new_y="NEXT", align="C")

        pdf.ln(6)
        pdf.set_draw_color(30, 58, 138)
        pdf.line(15, pdf.get_y(), 195, pdf.get_y())
        pdf.ln(8)

        pdf.set_text_color(30, 30, 30)
        for line in summary.split("\n"):
            line = line.strip()
            if not line:
                pdf.ln(3)
                continue

            if line.startswith("##"):
                pdf.ln(4)
                pdf.set_font("Helvetica", "B", 13)
                pdf.set_text_color(30, 58, 138)
                pdf.multi_cell(0, 7, line.replace("##", "").strip())
                pdf.ln(2)
            elif line.startswith(("•", "-", "*")):
                pdf.set_font("Helvetica", "", 10)
                pdf.set_text_color(40, 40, 40)
                clean = line.lstrip("•-* ").strip()
                pdf.multi_cell(0, 6, f"  -  {clean}")
                pdf.ln(1)
            else:
                pdf.set_font("Helvetica", "", 10)
                pdf.set_text_color(40, 40, 40)
                pdf.multi_cell(0, 6, line)
                pdf.ln(1)

        file_path = os.path.join(EXPORTS_DIR, f"{task_id}.pdf")
        pdf.output(file_path)

        task.status = "completed"
        task.file_url = f"/static/exports/{task_id}.pdf"
    except Exception as e:
        task.status = "failed"
        task.error_message = _sanitize_error_message(e)
    finally:
        task.updated_at = datetime.utcnow()
        db.commit()
        db.close()


# ── Background Task: PPTX Slides Generation ─────────────────────────────────
def generate_ppt_task(raw_text: str, task_id: str):
    db = SessionLocal()
    task = db.query(GenerationTask).filter(GenerationTask.task_id == task_id).first()
    if not task:
        db.close()
        return

    try:
        # Load API keys
        provider = os.getenv("AI_PROVIDER", "gemini").lower().strip()
        api_key = os.getenv("GOOGLE_API_KEY" if provider == "gemini" else "OPENAI_API_KEY", "")
        llm = get_llm(provider, api_key)

        prompt = """You are an expert presentation designer. Analyze the following document and 
extract exactly 4 key themes. For each theme, create a slide outline.

Return your response as a valid JSON array with exactly 4 objects. Each object must have:
- "title": A concise slide title (max 8 words)
- "bullets": An array of 3-5 concise bullet point strings
- "keyword": A single search keyword or short phrase (1-2 words) that describes the visual theme of the slide (e.g., "robotics", "brain", "calculation", "history", "graph", "business", "coding", "physics", "drone")

Document text:
{text}"""

        truncated = raw_text[:15000] if len(raw_text) > 15000 else raw_text
        response = llm.invoke([HumanMessage(content=prompt.format(text=truncated))])
        
        # Clean slide json
        cleaned_json = response.content.strip()
        cleaned_json = re.sub(r"^```(?:json)?\s*", "", cleaned_json)
        cleaned_json = re.sub(r"\s*```$", "", cleaned_json)
        match = re.search(r"\[.*\]", cleaned_json, re.DOTALL)
        if match:
            slides_data = json.loads(match.group())
        else:
            slides_data = json.loads(cleaned_json)

        # Build Presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Title Slide
        title_layout = prs.slide_layouts[6]
        title_slide = prs.slides.add_slide(title_layout)
        bg = title_slide.shapes.add_shape(1, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(30, 41, 59)
        bg.line.fill.background()

        tb = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "📚 AI Study Companion"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        sub = tf.add_paragraph()
        sub.text = "Key Themes & Visual Study Guide"
        sub.font.size = Pt(18)
        sub.font.color.rgb = RGBColor(148, 163, 184)
        sub.alignment = PP_ALIGN.CENTER

        # Content Slides
        for i, slide_data in enumerate(slides_data[:4]):
            title = slide_data.get("title", f"Theme {i + 1}")
            bullets = slide_data.get("bullets", ["No content generated."])
            keyword = slide_data.get("keyword", "education")

            slide = prs.slides.add_slide(title_layout)
            
            # Left accent
            accent_colors = [RGBColor(59, 130, 246), RGBColor(16, 185, 129), RGBColor(139, 92, 246), RGBColor(245, 158, 11)]
            color = accent_colors[i % len(accent_colors)]
            accent = slide.shapes.add_shape(1, Emu(0), Emu(0), Inches(0.15), prs.slide_height)
            accent.fill.solid()
            accent.fill.fore_color.rgb = color
            accent.line.fill.background()

            # Slide number badge
            badge = slide.shapes.add_shape(1, Inches(0.5), Inches(0.4), Inches(0.5), Inches(0.5))
            badge.fill.solid()
            badge.fill.fore_color.rgb = color
            badge.line.fill.background()
            badge.text_frame.text = str(i + 1)
            badge.text_frame.paragraphs[0].font.size = Pt(18)
            badge.text_frame.paragraphs[0].font.bold = True
            badge.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Title
            tbox = slide.shapes.add_textbox(Inches(1.3), Inches(0.35), Inches(7.5), Inches(0.7))
            tbox.text_frame.word_wrap = True
            tp = tbox.text_frame.paragraphs[0]
            tp.text = title
            tp.font.size = Pt(26)
            tp.font.bold = True
            tp.font.color.rgb = RGBColor(30, 41, 59)

            # Divider line
            line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(8.5), Inches(0.02))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(226, 232, 240)
            line.line.fill.background()

            # Download Image
            img_bytes = None
            if keyword and keyword.strip():
                try:
                    img_url = _get_exact_image_url(keyword)
                    res = requests.get(img_url, timeout=5)
                    if res.status_code == 200:
                        img_bytes = res.content
                except Exception:
                    pass

            box_width = Inches(4.7) if img_bytes else Inches(8.0)

            # Bullet box
            bbox = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), box_width, Inches(4.5))
            tf_bullets = bbox.text_frame
            tf_bullets.word_wrap = True

            for j, bullet in enumerate(bullets):
                bp = tf_bullets.paragraphs[0] if j == 0 else tf_bullets.add_paragraph()
                bp.text = f"▸  {bullet}"
                bp.font.size = Pt(16)
                bp.font.color.rgb = RGBColor(51, 65, 85)
                bp.space_after = Pt(12)

            if img_bytes:
                try:
                    slide.shapes.add_picture(
                        io.BytesIO(img_bytes),
                        Inches(5.5),
                        Inches(1.7),
                        width=Inches(3.8),
                        height=Inches(4.5)
                    )
                except Exception:
                    pass

        file_path = os.path.join(EXPORTS_DIR, f"{task_id}.pptx")
        prs.save(file_path)

        task.status = "completed"
        task.file_url = f"/static/exports/{task_id}.pptx"
    except Exception as e:
        task.status = "failed"
        task.error_message = _sanitize_error_message(e)
    finally:
        task.updated_at = datetime.utcnow()
        db.commit()
        db.close()


# ── Background Task: Audio Narration Generation ─────────────────────────────
def generate_audio_task(raw_text: str, task_id: str):
    db = SessionLocal()
    task = db.query(GenerationTask).filter(GenerationTask.task_id == task_id).first()
    if not task:
        db.close()
        return

    try:
        # Load API keys
        provider = os.getenv("AI_PROVIDER", "gemini").lower().strip()
        api_key = os.getenv("GOOGLE_API_KEY" if provider == "gemini" else "OPENAI_API_KEY", "")
        llm = get_llm(provider, api_key)

        prompt = """You are a friendly and expert study narrator. Read the following document in its entirety 
and create a highly detailed, comprehensive conversational review script that systematically explains 
the entire file and its contents, including all major concepts, sections, and findings.

Requirements:
- Write a thorough and rich spoken review script (4-6 detailed paragraphs, roughly 350-500 words).
- Begin with a friendly introduction explaining what the document is about overall.
- Walk through each of the main sections, themes, or core concepts in detail.
- Summarize the key takeaways and final conclusions at the end.
- Use a warm, engaging, educational, and conversational tone (as if teaching a student).
- Avoid jargon, special characters, bullet points, asterisks, or markdown formatting — output ONLY plain spoken paragraphs.

Document text:
{text}

Comprehensive Audio Script:"""

        truncated = raw_text[:40000] if len(raw_text) > 40000 else raw_text
        response = llm.invoke([HumanMessage(content=prompt.format(text=truncated))])
        script = response.content.strip()

        # Convert to MP3
        tts = gTTS(text=script, lang="en", slow=False)
        file_path = os.path.join(EXPORTS_DIR, f"{task_id}.mp3")
        tts.save(file_path)

        task.status = "completed"
        task.file_url = f"/static/exports/{task_id}.mp3"
    except Exception as e:
        task.status = "failed"
        task.error_message = _sanitize_error_message(e)
    finally:
        task.updated_at = datetime.utcnow()
        db.commit()
        db.close()


# ── Background Task: Flashcard Generation ───────────────────────────────────
def generate_flashcards_task(raw_text: str, session_id: str, task_id: str):
    db = SessionLocal()
    task = db.query(GenerationTask).filter(GenerationTask.task_id == task_id).first()
    if not task:
        db.close()
        return

    try:
        # Load API keys
        provider = os.getenv("AI_PROVIDER", "gemini").lower().strip()
        api_key = os.getenv("GOOGLE_API_KEY" if provider == "gemini" else "OPENAI_API_KEY", "")
        llm = get_llm(provider, api_key)

        prompt = """You are an expert educator. Read the following document and create 
exactly 5 high-quality study flashcards that test understanding of the key concepts.

Return your response as a valid JSON array with exactly 5 objects. Each object must have:
- "question": A clear, specific question testing a key concept
- "answer": A concise but complete answer (1-3 sentences)

IMPORTANT: Return ONLY the JSON array. No markdown, no code fences, no extra text.

Document text:
{text}"""

        truncated = raw_text[:15000] if len(raw_text) > 15000 else raw_text
        response = llm.invoke([HumanMessage(content=prompt.format(text=truncated))])

        cleaned = response.content.strip()
        cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
        cleaned = re.sub(r"\s*```$", "", cleaned)
        match = re.search(r"\[.*\]", cleaned, re.DOTALL)
        if match:
            cards = json.loads(match.group())
        else:
            cards = json.loads(cleaned)

        # Remove any existing flashcards for this session first
        db.query(Flashcard).filter(Flashcard.session_id == session_id).delete()

        # Save flashcards to SQLite database
        for card_data in cards[:5]:
            q = card_data.get("question", "No question generated.")
            a = card_data.get("answer", "No answer generated.")
            flashcard = Flashcard(session_id=session_id, question=q, answer=a)
            db.add(flashcard)

        task.status = "completed"
        task.file_url = f"/api/session/{session_id}/flashcards"  # Pointer to where flashcards can be loaded
    except Exception as e:
        task.status = "failed"
        task.error_message = _sanitize_error_message(e)
    finally:
        task.updated_at = datetime.utcnow()
        db.commit()
        db.close()
